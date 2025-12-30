from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import re
# from bs4 import BeautifulSoup

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors
DARK_BG = RGBColor(0x12, 0x12, 0x12)
ACCENT_BLUE = RGBColor(0x8B, 0xAF, 0xD0)
WHITE = RGBColor(0xE8, 0xE8, 0xE8)
GRAY = RGBColor(0xB3, 0xB3, 0xB3)
CARD_BG = RGBColor(0x1E, 0x1E, 0x1E)

def add_dark_background(slide):
    """Add dark background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_slide(title, subtitle, contact_info=None):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if contact_info:
        contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(2))
        tf = contact_box.text_frame
        for line in contact_info:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_items, is_grid=False):
    """Add a content slide with bullet points or grid"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(14.5), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if is_grid and len(content_items) == 3:
        # 3-column grid
        for i, (item_title, item_desc) in enumerate(content_items):
            x = Inches(0.75 + i * 5)
            
            # Card background
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), Inches(4.5), Inches(5.5))
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
            
            # Card title
            card_title = slide.shapes.add_textbox(x + Inches(0.25), Inches(2.5), Inches(4), Inches(1))
            tf = card_title.text_frame
            p = tf.paragraphs[0]
            p.text = item_title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # Card description
            card_desc = slide.shapes.add_textbox(x + Inches(0.25), Inches(3.5), Inches(4), Inches(3.5))
            tf = card_desc.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item_desc
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER
    else:
        # Bullet points
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(14.5), Inches(6.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(24)
            p.font.color.rgb = WHITE
            p.space_after = Pt(16)

def add_metrics_slide(title, metrics):
    """Add a slide with key metrics"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(14.5), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Metrics in grid
    cols = min(len(metrics), 4)
    col_width = 14.5 / cols
    
    for i, (value, label) in enumerate(metrics):
        row = i // 4
        col = i % 4
        x = Inches(0.75 + col * col_width)
        y = Inches(2.5 + row * 3)
        
        # Value
        val_box = slide.shapes.add_textbox(x, y, Inches(col_width - 0.5), Inches(1.5))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(x, y + Inches(1.2), Inches(col_width - 0.5), Inches(1))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

# Slide 1: Title
add_title_slide(
    "ArdiA Health Labs",
    "FROM ALERTING TO SOLVING",
    ["Ram Vadlamudi | Founder", "ram.vadlamudi@ardiahealthlabs.com | +1(469)679-3334"]
)

# Slide 2: The Problem
add_content_slide("The Problem", [
    "Chronic diseases cost $4.1 trillion annually in the US",
    "80% of healthcare costs come from 20% of patients with chronic conditions",
    "Current remote monitoring systems alert but don't solve",
    "Providers are overwhelmed with data but lack actionable insights",
    "Patients experience fragmented, reactive care"
])

# Slide 3: Market Opportunity
add_metrics_slide("Market Opportunity", [
    ("$50B", "Remote Patient Monitoring Market by 2027"),
    ("133M", "Americans with Chronic Conditions"),
    ("25%", "CAGR in Digital Health"),
    ("$4.1T", "Annual Chronic Disease Costs")
])

# Slide 4: Our Solution
add_content_slide("Our Solution: Clinical Reasoning Engine", [
    "AI-powered platform that goes beyond alerting to SOLVING",
    "Integrates patient data, environmental factors, and clinical patterns",
    "Predicts exacerbations 72+ hours before symptoms appear",
    "Delivers personalized, actionable interventions to patients and providers",
    "Seamless EHR integration with existing workflows"
])

# Slide 5: How It Works
add_content_slide("How It Works", [
    "COLLECT: Continuous monitoring of vitals, symptoms, and environmental data",
    "ANALYZE: AI correlates patterns with historical episodes and clinical guidelines",
    "PREDICT: Risk scoring identifies patients heading toward acute episodes",
    "INTERVENE: Automated, personalized recommendations to patients and provider alerts",
    "LEARN: System continuously improves based on outcomes"
])

# Slide 6: Technology Platform
add_content_slide("Technology Platform", [
    "Multi-modal data integration (wearables, EHR, environmental, voice biomarkers)",
    "Advanced reasoning engine with clinical knowledge graphs",
    "Real-time risk stratification and prediction models",
    "Natural language patient communication",
    "HIPAA-compliant cloud infrastructure"
])

# Slide 7: Three Stakeholder Value
add_content_slide("Value for Every Stakeholder", [
    ("Patients", "Better outcomes, fewer ER visits, and personalized 24/7 care"),
    ("Providers", "Clinical efficiency, reduced burnout, and new revenue streams"),
    ("Payers", "Significant cost savings, risk stratification, and quality improvements")
], is_grid=True)

# Slide 8: Patient Value
add_metrics_slide("Patient Value: Better Health Outcomes", [
    ("-74%", "ER Visit Reduction"),
    ("0", "Hospitalizations in Pilot"),
    ("94%", "Patient Engagement Rate"),
    ("4.8/5", "Patient Satisfaction")
])

# Slide 9: Provider Value - Efficiency
add_metrics_slide("Provider Value: Efficiency & Revenue", [
    ("-68%", "Reduced Admin Time"),
    ("4.2 hrs", "Saved Per Week/Provider"),
    ("$22K", "Projected Annual Revenue/Clinic"),
    ("96%", "Alert Acceptance Rate")
])

# Slide 10: Provider Testimonials
add_content_slide("Provider Testimonials", [
    '"The Clinical Reasoning Engine doesn\'t just alert me—it tells me WHY. This is the future of proactive care." — Dr. Sarah Mitchell, MD, Pulmonology',
    '"Dorothy\'s fluid status is now managed proactively. The engine caught a problem 3 days before she would have shown up in my ER." — Dr. James Rodriguez, MD, FASN, Nephrology'
])

# Slide 11: Insurance Value - Financial
add_metrics_slide("Payer Value: Financial Impact (per 1,000 members)", [
    ("$8.1M", "Total Annual Savings"),
    ("$2.4M", "ER Cost Avoidance"),
    ("$4.8M", "Hospitalization Savings"),
    ("$890K", "Medication Optimization")
])

# Slide 12: Insurance Value - Quality
add_metrics_slide("Payer Value: Quality & Risk", [
    ("94%", "Prediction Accuracy"),
    ("72 hrs", "Early Detection Window"),
    ("+18%", "Asthma Med Ratio (HEDIS)"),
    ("4→5", "Star Rating Potential")
])

# Slide 13: Pilot Results
add_metrics_slide("Pilot Program Results (90 Days)", [
    ("3", "Active Patients"),
    ("2", "Partner Clinics"),
    ("$29,790", "Cost Savings Achieved"),
    ("47:1", "ROI Ratio")
])

# Slide 14: Business Model
add_content_slide("Business Model", [
    "B2B2C SaaS model with multiple revenue streams:",
    "• Clinic Subscription: $500-2,000/month per practice",
    "• Per Patient Fee: $50-150/month for chronic care management",
    "• Payer Contracts: Risk-based, PMPM, or value-based arrangements",
    "• CCM/RTM Billing Support: Enable practices to capture Medicare revenue"
])

# Slide 15: Roadmap
add_content_slide("Roadmap", [
    "Q1 2025: Expand pilot to 10 clinics, 50 patients",
    "Q2 2025: Launch payer partnership pilots",
    "Q3 2025: Series A funding, expand to 3 chronic conditions",
    "Q4 2025: 100 clinic partnerships",
    "2026: National expansion, 5 chronic conditions, payer integrations"
])

# Slide 16: The Ask
add_metrics_slide("The Ask", [
    ("$2M", "Seed Round"),
    ("18 mo", "Runway"),
    ("10x", "Clinic Expansion"),
    ("$50K", "Monthly Revenue Target")
])

# Slide 17: Contact
add_title_slide(
    "Invest in the Future",
    '"Let\'s solve healthcare together"',
    ["Ram Vadlamudi, Founder", "ram.vadlamudi@ardiahealthlabs.com", "+1 (469) 679-3334", "ardiahealthlabs.com"]
)

# Save
output_path = "/workspace/ArdiA_Health_Labs_Investor_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
